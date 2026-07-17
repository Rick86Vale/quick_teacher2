from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re

def criar_apresentacao(titulo_aula, conteudo_markdown):
    prs = Presentation()
    
    # Identidade Visual (Verde do seu CSS)
    COR_PRINCIPAL = RGBColor(22, 172, 117)
    
    # 1. Limpeza de caracteres de controle e quebras de linha Windows/Linux
    conteudo = conteudo_markdown.replace('\r\n', '\n').replace('\r', '\n')
    
    # 2. Divisão por --- (quebra manual de slide) ou por títulos (quebra automática)
    # Aqui usamos --- como divisor principal de slides
    secoes = re.split(r'\n---\n', conteudo)
    
    # Slide de Título (Capa)
    slide_capa = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide_capa.shapes.title
    subtitle = slide_capa.shapes.placeholders[1]
    title.text = titulo_aula
    title.text_frame.paragraphs[0].font.color.rgb = COR_PRINCIPAL
    subtitle.text = "Material de apoio gerado automaticamente"

    # 3. Processamento de cada slide (O LOOP COMEÇA AQUI)
    for secao in secoes:
        secao = secao.strip()
        if not secao:
            continue
            
        # --- A LÓGICA DE EXTRAÇÃO ESTÁ AQUI DENTRO ---
        linhas = secao.split('\n')
        
        # Se a primeira linha começa com #, ela é o título e deve ser removida do corpo
        if linhas[0].strip().startswith('#'):
            titulo_slide = linhas[0].replace('#', '').strip()
            # O conteúdo será tudo o que estiver abaixo da primeira linha
            conteudo_slide = '\n'.join(linhas[1:]).strip()
        else:
            # Caso não haja título, usamos um título genérico
            titulo_slide = "Tópico da Aula"
            conteudo_slide = secao.strip()
        
        # Limpeza final dos caracteres de controle
        conteudo_slide = conteudo_slide.replace('\r', '')

        # Criação do Slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Estilização do Título
        title_shape = slide.shapes.title
        title_shape.text = titulo_slide
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = COR_PRINCIPAL
            paragraph.font.bold = True
            paragraph.font.size = Pt(36)
            
        # Estilização do Conteúdo
        body_shape = slide.shapes.placeholders[1]
        body_shape.text = conteudo_slide
        for paragraph in body_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = RGBColor(50, 50, 50)

    return prs